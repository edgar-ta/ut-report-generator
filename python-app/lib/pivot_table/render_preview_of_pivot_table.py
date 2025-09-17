from control_variables import PATH_OF_PPTX_TEMPLATE

from lib.directory_definitions import temporary_rendered_file_of_report, preview_image_of_slide

from models.pivot_table.self import PivotTable
from models.report.self import Report

from render.drawable_area import DrawableArea

from pptx import Presentation
from pptx.slide import Slide
from pptx.util import Cm

from spire.presentation import Presentation as SpirePresentation

import os

def render_preview_of_pivot_table(pivot_table: PivotTable, root_directory: str):
    template_path = PATH_OF_PPTX_TEMPLATE()
    temporary_path = temporary_rendered_file_of_report(root_directory=root_directory)

    presentation = Presentation(template_path)
    
    base_area = DrawableArea(
        x=0, y=0, 
        width=presentation.slide_width.emu, 
        height=presentation.slide_height.emu
        ).with_padding(horizontal=Cm(2).emu, vertical=Cm(4).emu)

    def fresh_slide() -> Slide:
        slide = presentation.slides.add_slide(presentation.slides[0].slide_layout)
        for placeholder in slide.placeholders:
            slide.placeholders.element.remove(placeholder.element)
        return slide

    pivot_table.render(slide=fresh_slide(), drawable_area=base_area)
    presentation.save(temporary_path)
    spire_presentation = SpirePresentation()
    spire_presentation.LoadFromFile(temporary_path)

    filepath = preview_image_of_slide(root_directory=root_directory, slide_id=pivot_table.identifier)

    spire_slide = spire_presentation.Slides[1]
    image = spire_slide.SaveAsImage()
    image.Save(filepath)
    image.Dispose()

    if pivot_table.preview is not None and os.path.exists(pivot_table.preview):
        os.remove(pivot_table.preview)
    pivot_table.preview = filepath

    spire_presentation.Dispose()
    os.remove(temporary_path)
