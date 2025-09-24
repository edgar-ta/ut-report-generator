from control_variables import PATH_OF_PPTX_TEMPLATE

from models.slide.self import Slide as ProjectSlide

from render.drawable_area import DrawableArea

from pptx import Presentation
from pptx.slide import Slide
from pptx.util import Cm

def compile_slides(slides: list[ProjectSlide], filepath: str):
    template_path = PATH_OF_PPTX_TEMPLATE()
    presentation = Presentation(template_path)
    
    base_area = DrawableArea(
        x=0, y=0, 
        width=presentation.slide_width.emu, 
        height=presentation.slide_height.emu
        ).with_padding(horizontal=Cm(2).emu, vertical=Cm(4).emu)

    def fresh_slide() -> Slide:
        slide = presentation.slides.add_slide(presentation.slides[0].slide_layout)
        for placeholder in slide.placeholders:
            slide.placeholders.element.remove(placeholder.element)
        return slide
    
    for slide in slides:
        pptx_slide = fresh_slide()
        slide.render(slide=pptx_slide, drawable_area=base_area)
    
    presentation.save(filepath)
