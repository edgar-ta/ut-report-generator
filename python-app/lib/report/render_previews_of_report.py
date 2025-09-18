from control_variables import PATH_OF_PPTX_TEMPLATE

from lib.directory_definitions import temporary_rendered_file_of_report, preview_image_of_slide

from models.report.self import Report

from render.drawable_area import DrawableArea

from pptx import Presentation
from pptx.slide import Slide
from pptx.util import Cm

from spire.presentation import Presentation as SpirePresentation

import os

def render_previews_of_report(report: Report):
    template_path = PATH_OF_PPTX_TEMPLATE()
    temporary_path = temporary_rendered_file_of_report(root_directory=report.root_directory)

    presentation = Presentation(template_path)
    
    base_area = DrawableArea(
        x=0, y=0, 
        width=presentation.slide_width.emu, 
        height=presentation.slide_height.emu
        ).with_padding(horizontal=Cm(2).emu, vertical=Cm(4).emu)

    def fresh_slide() -> Slide:
        slide = presentation.slides.add_slide(presentation.slides[0].slide_layout)
        for placeholder in slide.placeholders:
            slide.placeholders.element.remove(placeholder.element)
        return slide
    
    for slide in report.slides:
        pptx_slide = fresh_slide()
        slide.render(slide=pptx_slide, drawable_area=base_area)
    
    presentation.save(temporary_path)
    spire_presentation = SpirePresentation()
    spire_presentation.LoadFromFile(temporary_path)

    for index, slide in enumerate(report.slides):
        spire_slide = spire_presentation.Slides[index + 1]
        file_name = preview_image_of_slide(root_directory=report.root_directory, slide_id=slide.identifier)

        image = spire_slide.SaveAsImage()
        image.Save(file_name)
        image.Dispose()

        slide.preview = file_name

    spire_presentation.Dispose()
    os.remove(temporary_path)
