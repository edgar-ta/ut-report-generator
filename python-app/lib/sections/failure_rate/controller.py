from lib.slide_controller import SlideController
from lib.descriptive_error import DescriptiveError
from lib.sections.failure_rate.source import read_excel, get_clean_data_frame, create_unit_name, graph_failure_rate
from lib.get_asset import get_string_asset, get_image_asset
from lib.get_or_panic import get_or_panic

from models.slide_type import SlideType
from models.asset import Asset

from pandas import DataFrame
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

import re
import os
import uuid

def delayed_teachers_legend(subjects_without_grades: DataFrame, unit: int) -> str:
    # Get unique professor names from the multi-index
    subjects = subjects_without_grades.index.get_level_values("subject").unique()
    subjects = [str(p).title() for p in subjects]

    if not subjects:
        return f"Todos los profesores subieron calificaciones en la unidad {unit}"

    if len(subjects) == 1:
        return f"La materia de {subjects[0]} no subió calificaciones en la unidad {unit}"
    elif len(subjects) == 2:
        return f"Las materias de {subjects[0]} y {subjects[1]} no subieron calificaciones en la unidad {unit}"
    else:
        *rest, last = subjects
        return f"Las materias de {', '.join(rest)}, y {last} no subieron calificaciones en la unidad {unit}"

class FailureRate_Controller(SlideController):
    @staticmethod
    def slide_type() -> SlideType:
        return SlideType.FAILURE_RATE

    @staticmethod
    def default_arguments() -> dict[str, any]:
        return {
            "unit": 1,
            "show_delayed_teachers": True
        }

    @staticmethod
    def build_assets(data_files: list[str], base_directory: str, arguments: dict[str, any]) -> list[Asset]:
        data_frame = read_excel(data_files[0])
        data_frame = get_clean_data_frame(data_frame)

        unit = arguments["unit"]
        grades = data_frame.xs(create_unit_name(unit), level="unit").map(lambda value: abs(value))

        subjects_without_grades = grades[(grades == 0).all(axis=1)]
        subjects_with_grades = grades[(grades != 0).any(axis=1)]

        main_chart_path = os.path.join(base_directory, str(uuid.uuid4()) + ".png")
        graph_failure_rate(subjects_with_grades, main_chart_path)

        return Asset.list_from([ 
            ("main_chart", main_chart_path, "image"), 
            ("delayed_teachers", delayed_teachers_legend(subjects_without_grades, unit), "text") 
        ])


    @staticmethod
    def validate_arguments(arguments: dict[str, any]) -> None:
        # UNIT
        unit = get_or_panic(arguments, "unit", "La unidad del gráfico debe estar presente")
        if unit > 5:
            raise DescriptiveError(400, "No subject has more than 5 units")

        # SHOW_DELAYED_TEACHERS
        show_delayed_teachers = get_or_panic(arguments, "show_delayed_teachers", "La opción de mostrar profesores con calificaciones atrasadas debe estar presente")
        if not type(show_delayed_teachers) is bool:
            raise DescriptiveError(400, f"Invalid value for 'show_delayed_teachers'. Expected 'true' or 'false', found {show_delayed_teachers}")


    @staticmethod
    def render_slide(presentation: Presentation, arguments: dict[str, any], assets: list[Asset]) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])

        main_chart_path = get_image_asset(assets, "main_chart")

        left = Inches(1)
        top = Inches(1)
        width = Inches(8)
        slide.shapes.add_picture(main_chart_path, left, top, width=width)

        if arguments["show_delayed_teachers"]:
            delayed_teachers_text = get_string_asset(assets, "delayed_teachers")

            left = Inches(1)
            top = Inches(6)
            width = Inches(8)
            height = Inches(1)
            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            text_frame.text = delayed_teachers_text

            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.bold = True
                paragraph.alignment = PP_ALIGN.CENTER
