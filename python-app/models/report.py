from lib.get_metadata import get_metadata
from lib.descriptive_error import DescriptiveError
from lib.directory_definitions import metadata_file_of_report, slides_directory_of_report, root_directory_of_report, rendered_file_of_report, export_file_of_report, export_directory_of_report, data_directory_of_report, get_reports_directory

from models.image_slide.self import ImageSlide
from models.pivot_table.self import PivotTable
from models.slide.slide_category import SlideCategory

from control_variables import CURRENT_DIRECTORY_PATH, CURRENT_PROJECT_VERSION

from pptx import Presentation
from functools import cached_property

import pandas
import os
import uuid
import json

Slide = ImageSlide | PivotTable

class Report:
    def __init__(
            self, 
            identifier: str,
            root_directory: str, 
            report_name: str, 
            creation_date: pandas.Timestamp, 
            last_edit: pandas.Timestamp,
            last_open: pandas.Timestamp,
            slides: list[ImageSlide | PivotTable],
            version: str
            ) -> None:
        # Data that goes to the serialization dict
        self.identifier = identifier
        self.report_name = report_name
        self.creation_date = creation_date
        self.last_edit = last_edit
        self.last_open = last_open
        self.slides = slides
        self.version = version

        # Data that does not go to the serialization dict
        self.root_directory = root_directory

    def to_dict(self) -> dict:
        """
        Serializes the Report instance to a dictionary for JSON export.
        Dates are converted to ISO 8601 strings.
        """
        return {
            "identifier": self.identifier,
            "report_name": self.report_name,
            "creation_date": self.creation_date.isoformat(),
            "last_edit": self.last_edit.isoformat(),
            "last_open": self.last_open.isoformat(),
            "slides": [slide.to_dict() for slide in self.slides],
            "version": self.version
        }
    
    @classmethod
    def from_root_directory(cls, root_directory: str) -> "Report":
        metadata = get_metadata(root_directory=root_directory)
        
        version = metadata['version']
        if version != CURRENT_PROJECT_VERSION:
            raise DescriptiveError(message=f"El reporte no se puede abrir porque es de una versión desactualizada. La versión actual es {CURRENT_PROJECT_VERSION} y el reporte tiene {version}", http_error_code=400)

        return cls(
            identifier=metadata['identifier'],
            root_directory=root_directory,
            report_name=metadata["report_name"],
            creation_date=pandas.Timestamp(metadata["creation_date"]),
            last_edit=pandas.Timestamp(metadata["last_edit"]),
            last_open=pandas.Timestamp(metadata["last_open"]),
            slides=[
                    ImageSlide.from_json(json_data=slide) 
                    if SlideCategory(slide["category"]) == SlideCategory.IMAGE_SLIDE
                    else PivotTable.from_json(json_data=slide)
                for slide in metadata.get("slides", [])
            ],
            version=version
        )

    @classmethod
    def from_identifier(cls, identifier: str) -> "Report":
        for filename in os.listdir(get_reports_directory()):
            directory_path = os.path.join(get_reports_directory(), filename)
            if os.path.isdir(directory_path) and filename.endswith(identifier):
                return Report.from_root_directory(root_directory=directory_path)
            
        raise DescriptiveError(http_error_code=400, message=f'La id especificada ({identifier}) no corresponde a ningún reporte')

    @classmethod
    def from_nothing(cls) -> "Report":
        report_id = cls.new_report_id()

        report = Report(
            identifier=report_id,
            root_directory=root_directory_of_report(report_id),
            report_name="Mi reporte",
            creation_date=pandas.Timestamp.now(),
            last_edit=pandas.Timestamp.now(),
            last_open=pandas.Timestamp.now(),
            slides=[],
            version=CURRENT_PROJECT_VERSION
        )
        return report
    
    def makedirs(self, exist_ok: bool = True) -> None:
        '''
        Creates the directories that the reports needs in order
        to work (i. e., the root, the slides and data directory)
        '''
        os.makedirs(self.root_directory, exist_ok=exist_ok)
        os.makedirs(self.slides_directory, exist_ok=exist_ok)
        os.makedirs(self.data_directory, exist_ok=exist_ok)

        for slide in self.slides:
            slide.makedirs(exist_ok=exist_ok)

    def new_render(self) -> None:
        is_render_updated = self.last_render is not None and self.last_render >= self.last_edit

        if is_render_updated and os.path.exists(self.rendered_file):
            return
        
        presentation = Presentation()
        for slide in self.slides:
            slide.build_new_assets()
            slide.controller.render_slide(
                presentation=presentation, 
                arguments=slide.arguments, 
                assets=slide.assets
            )
        
        if os.path.exists(self.rendered_file):
            os.remove(self.rendered_file)
        
        presentation.save(self.rendered_file)

    def add_slide(self, slide: Slide) -> None:
        self.slides.append(slide)
    
    def save(self):
        last_edits = [ slide.last_edit for slide in self.slides ]
        last_edits.append(self.last_edit)
        self.last_edit = max(last_edits)

        with open(self.metadata_file, "w") as json_file:
            json.dump(self.to_dict(), json_file, indent=4)

    def __getitem__(self, key: str) -> Slide:
        slide = next((slide for slide in self.slides if slide.identifier == key), None)
        if slide is None:
            raise DescriptiveError(400, f"La diapositiva con id {key} no existe. Tal vez sea un error de dedo")
        return slide

    @classmethod
    def get_reports_directory(cls) -> str:
        return os.path.join(CURRENT_DIRECTORY_PATH, "reports")

    @cached_property
    def slides_directory(self) -> str:
        return slides_directory_of_report(self.root_directory)
    
    @cached_property
    def metadata_file(self) -> str:
        return metadata_file_of_report(self.root_directory)
    
    @property
    def rendered_file(self) -> str:
        return rendered_file_of_report(root_directory=self.root_directory, report_name=self.report_name)
    
    @property
    def export_file(self) -> str:
        return export_file_of_report(root_directory=self.root_directory, report_name=self.report_name)
    
    @cached_property
    def export_directory(self) -> str:
        return export_directory_of_report(root_directory=self.root_directory)
    
    @cached_property
    def data_directory(self) -> str:
        return data_directory_of_report(root_directory=self.root_directory)
    
    @classmethod
    def new_report_id(cls) -> str:
        return str(uuid.uuid4())
