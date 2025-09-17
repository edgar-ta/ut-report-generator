from render.widget.drawable_widget import DrawableWidget

from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN
from render.font import Font
from render.color import Color

class TextWidget(DrawableWidget):
    def __init__(
            self, 
            text: str, 
            alignment: PP_ALIGN = PP_ALIGN.LEFT,
            font: Font = Font(),
            color: Color = Color(0, 0, 0)
            ):
        super().__init__()

        self.text = text
        self.alignment = alignment
        self.font = font
        self.color = color

    def draw(self, slide, drawable_area):
        textbox = slide.shapes.add_textbox(
            left=drawable_area.x, 
            top=drawable_area.y, 
            width=Emu(drawable_area.width), 
            height=Emu(drawable_area.height)
            )
        
        text_frame = textbox.text_frame
        text_frame.clear()

        for line in self.text.splitlines():
            paragraph = text_frame.add_paragraph()
            paragraph.text = line
            paragraph.alignment = self.alignment
            text_frame.word_wrap = True

            paragraph.font.color.rgb = self.color
            paragraph.font.name = self.font.font_family
            paragraph.font.size = self.font.size
            paragraph.font.bold = self.font.bold
            paragraph.font.italic = self.font.italic
