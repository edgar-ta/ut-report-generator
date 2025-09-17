from routes.hello_world import hello_world

import routes.pivot_table as pivot_table
import routes.image_slide as image_slide
import routes.report as report

from flask import Flask, request, jsonify
from typing import TypeVar
from dotenv import load_dotenv

import logging
import sys

app = Flask(__name__)
logger = logging.getLogger(__name__)

hello_world(app)

app.register_blueprint(report.blueprint)
app.register_blueprint(pivot_table.blueprint)
app.register_blueprint(image_slide.blueprint)

T = TypeVar("T")
def item_or(_list: list[T], index: int, default: T) -> T:
    try:
        return _list[index]
    except IndexError:
        return default

if __name__ == '__main__':
    mode = item_or(_list=sys.argv, index=1, default="dev")
    port = int(item_or(_list=sys.argv, index=2, default="5000"))
    load_dotenv()

    if mode == "dev":
        logging.basicConfig(filename='logs.log')
        app.run(port=port)
    elif mode == "release":
        from waitress import serve
        serve(app=app, port=port)
    elif mode == "playground":
        from models.image_slide.cover_page_slide import CoverPageSlide
        from pptx import Presentation
        from pptx.util import Cm
        from pandas import Timestamp

        from render.drawable_area import DrawableArea

        path = r'D:\college\cuatrimestre-6\2025-06-16--estadias\ut-report-generator\.logistics-assets\example-presentation-copy.pptx'

        presentation = Presentation(path)
        slide = presentation.slides.add_slide(presentation.slides[0].slide_layout)
        base_area = DrawableArea(
            x=0, y=0, 
            width=presentation.slide_width.emu, 
            height=presentation.slide_height.emu
            ).with_padding(horizontal=Cm(2).emu, vertical=Cm(4).emu)
        
        for placeholder in slide.placeholders:
            slide.placeholders.element.remove(placeholder.element)

        cover_page = CoverPageSlide(
            title="1ER. INFORME GRUPO DS01SM-24-2°",
            identifier="fsfs",
            creation_date=Timestamp.now(),
            last_edit=Timestamp.now(),
            preview=None,
            professor_name="BRENDA JUÁREZ SANTIAGO",
            period="ENERO - ABRIL 2025",
            date="16/09/25"
        )
        cover_page.render(slide=slide, drawable_area=base_area)
        
        presentation.save(path)
